"""Document processing service for extracting text and creating chunks."""

import os
import re
import uuid
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
import logging

# Document processing libraries
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import requests

from ..models.models import Document, DocumentChunk
from ..schemas.documents import DocumentCreate


logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Service for processing documents and extracting text content."""
    
    def __init__(self, chunk_size: int = 700, chunk_overlap: int = 100):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.supported_types = ["pdf", "docx", "pptx", "txt", "gdrive"]
    
    async def process_document(self, document: Document) -> bool:
        """Process a document and create chunks."""
        try:
            logger.info(f"Processing document: {document.name} ({document.file_type})")
            
            # Extract text based on file type
            text_content = await self._extract_text(document)
            
            if not text_content:
                logger.warning(f"No text extracted from document: {document.name}")
                return False
            
            # Create chunks
            chunks = self._create_chunks(text_content, document.id)
            
            # Save chunks to database (this would be done in the calling service)
            logger.info(f"Created {len(chunks)} chunks for document: {document.name}")
            
            return True
            
        except Exception as e:
            logger.error(f"Error processing document {document.name}: {str(e)}")
            return False
    
    async def _extract_text(self, document: Document) -> str:
        """Extract text content from document based on file type."""
        if not os.path.exists(document.file_path):
            raise FileNotFoundError(f"Document file not found: {document.file_path}")
        
        if document.file_type == "pdf":
            return self._extract_pdf_text(document.file_path)
        elif document.file_type == "docx":
            return self._extract_docx_text(document.file_path)
        elif document.file_type == "pptx":
            return self._extract_pptx_text(document.file_path)
        elif document.file_type == "txt":
            return self._extract_txt_text(document.file_path)
        elif document.file_type == "gdrive":
            return await self._extract_gdrive_text(document.file_path)
        else:
            raise ValueError(f"Unsupported file type: {document.file_type}")
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF file."""
        text_content = []
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            # Add page marker for citation purposes
                            text_content.append(f"[PAGE {page_num}]\n{page_text}")
                    except Exception as e:
                        logger.warning(f"Error extracting text from page {page_num}: {e}")
                        continue
                
        except Exception as e:
            logger.error(f"Error reading PDF file {file_path}: {e}")
            raise
        
        return "\n\n".join(text_content)
    
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            doc = DocxDocument(file_path)
            text_content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            logger.error(f"Error reading DOCX file {file_path}: {e}")
            raise
    
    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PPTX file."""
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"[SLIDE {slide_num}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the slide marker
                    text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            logger.error(f"Error reading PPTX file {file_path}: {e}")
            raise
    
    def _extract_txt_text(self, file_path: str) -> str:
        """Extract text from TXT file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Error reading TXT file {file_path}: {e}")
                raise
        except Exception as e:
            logger.error(f"Error reading TXT file {file_path}: {e}")
            raise
    
    async def _extract_gdrive_text(self, gdrive_url: str) -> str:
        """Extract text from Google Drive document."""
        # This is a placeholder implementation
        # In production, you would use Google Drive API
        logger.warning("Google Drive extraction not implemented yet")
        return f"Google Drive document: {gdrive_url}\n\nContent extraction from Google Drive requires API integration."
    
    def _create_chunks(self, text: str, document_id: str) -> List[DocumentChunk]:
        """Create text chunks from document content."""
        # Clean and normalize text
        cleaned_text = self._clean_text(text)
        
        # Split into sentences for better chunk boundaries
        sentences = self._split_into_sentences(cleaned_text)
        
        chunks = []
        current_chunk = ""
        current_tokens = 0
        chunk_index = 0
        
        for sentence in sentences:
            sentence_tokens = self._estimate_tokens(sentence)
            
            # If adding this sentence would exceed chunk size, create a new chunk
            if current_tokens + sentence_tokens > self.chunk_size and current_chunk:
                chunk = self._create_chunk(
                    content=current_chunk.strip(),
                    document_id=document_id,
                    chunk_index=chunk_index
                )
                chunks.append(chunk)
                
                # Start new chunk with overlap
                overlap_text = self._get_overlap_text(current_chunk, self.chunk_overlap)
                current_chunk = overlap_text + " " + sentence
                current_tokens = self._estimate_tokens(current_chunk)
                chunk_index += 1
            else:
                current_chunk += " " + sentence
                current_tokens += sentence_tokens
        
        # Add final chunk if there's remaining content
        if current_chunk.strip():
            chunk = self._create_chunk(
                content=current_chunk.strip(),
                document_id=document_id,
                chunk_index=chunk_index
            )
            chunks.append(chunk)
        
        return chunks
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text content."""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters that might interfere with processing
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)\[\]\{\}\"\'\/]', '', text)
        
        # Normalize line breaks
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        
        return text.strip()
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences for better chunk boundaries."""
        # Simple sentence splitting - could be improved with NLTK
        sentence_endings = r'[.!?]+\s+'
        sentences = re.split(sentence_endings, text)
        
        # Filter out empty sentences and very short ones
        sentences = [s.strip() for s in sentences if len(s.strip()) > 10]
        
        return sentences
    
    def _estimate_tokens(self, text: str) -> int:
        """Estimate token count for text (rough approximation)."""
        # Simple approximation: ~4 characters per token
        return len(text) // 4
    
    def _get_overlap_text(self, text: str, overlap_tokens: int) -> str:
        """Get overlap text from the end of current chunk."""
        words = text.split()
        # Approximate: take last N words where N * 4 ≈ overlap_tokens
        overlap_words = overlap_tokens // 4
        if len(words) > overlap_words:
            return " ".join(words[-overlap_words:])
        return text
    
    def _create_chunk(self, content: str, document_id: str, chunk_index: int) -> DocumentChunk:
        """Create a DocumentChunk object."""
        chunk_id = str(uuid.uuid4())
        
        # Extract page number if present in content
        page_number = self._extract_page_number(content)
        
        # Extract section if present
        section = self._extract_section(content)
        
        return DocumentChunk(
            id=chunk_id,
            document_id=document_id,
            content=content,
            page_number=page_number,
            section=section,
            token_count=self._estimate_tokens(content),
            chunk_index=chunk_index,
            created_at=datetime.utcnow()
        )
    
    def _extract_page_number(self, content: str) -> Optional[int]:
        """Extract page number from content if present."""
        page_match = re.search(r'\[PAGE (\d+)\]', content)
        if page_match:
            return int(page_match.group(1))
        return None
    
    def _extract_section(self, content: str) -> Optional[str]:
        """Extract section information from content."""
        # Look for slide markers
        slide_match = re.search(r'\[SLIDE (\d+)\]', content)
        if slide_match:
            return f"Slide {slide_match.group(1)}"
        
        # Look for common section headers
        section_patterns = [
            r'^(Chapter \d+)',
            r'^(Section \d+)',
            r'^([A-Z][A-Za-z\s]+):',
            r'^(\d+\.\s+[A-Za-z\s]+)'
        ]
        
        for pattern in section_patterns:
            match = re.search(pattern, content, re.MULTILINE)
            if match:
                return match.group(1).strip()
        
        return None
    
    def get_document_metadata(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """Extract metadata from document file."""
        metadata = {
            "file_size": os.path.getsize(file_path) if os.path.exists(file_path) else 0,
            "processed_at": datetime.utcnow().isoformat()
        }
        
        try:
            if file_type == "pdf":
                metadata.update(self._get_pdf_metadata(file_path))
            elif file_type == "docx":
                metadata.update(self._get_docx_metadata(file_path))
            elif file_type == "pptx":
                metadata.update(self._get_pptx_metadata(file_path))
        except Exception as e:
            logger.warning(f"Could not extract metadata from {file_path}: {e}")
        
        return metadata
    
    def _get_pdf_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get PDF-specific metadata."""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                metadata = {
                    "page_count": len(pdf_reader.pages)
                }
                
                if pdf_reader.metadata:
                    if pdf_reader.metadata.get('/Author'):
                        metadata["author"] = pdf_reader.metadata['/Author']
                    if pdf_reader.metadata.get('/Title'):
                        metadata["title"] = pdf_reader.metadata['/Title']
                    if pdf_reader.metadata.get('/CreationDate'):
                        metadata["creation_date"] = str(pdf_reader.metadata['/CreationDate'])
                
                return metadata
        except Exception:
            return {"page_count": 0}
    
    def _get_docx_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get DOCX-specific metadata."""
        try:
            doc = DocxDocument(file_path)
            metadata = {}
            
            if hasattr(doc.core_properties, 'author') and doc.core_properties.author:
                metadata["author"] = doc.core_properties.author
            if hasattr(doc.core_properties, 'title') and doc.core_properties.title:
                metadata["title"] = doc.core_properties.title
            if hasattr(doc.core_properties, 'created') and doc.core_properties.created:
                metadata["creation_date"] = doc.core_properties.created.isoformat()
            
            # Count paragraphs as rough page estimate
            paragraph_count = len([p for p in doc.paragraphs if p.text.strip()])
            metadata["paragraph_count"] = paragraph_count
            metadata["estimated_pages"] = max(1, paragraph_count // 20)  # Rough estimate
            
            return metadata
        except Exception:
            return {}
    
    def _get_pptx_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get PPTX-specific metadata."""
        try:
            prs = Presentation(file_path)
            metadata = {
                "slide_count": len(prs.slides)
            }
            
            if hasattr(prs.core_properties, 'author') and prs.core_properties.author:
                metadata["author"] = prs.core_properties.author
            if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
                metadata["title"] = prs.core_properties.title
            if hasattr(prs.core_properties, 'created') and prs.core_properties.created:
                metadata["creation_date"] = prs.core_properties.created.isoformat()
            
            return metadata
        except Exception:
            return {"slide_count": 0}